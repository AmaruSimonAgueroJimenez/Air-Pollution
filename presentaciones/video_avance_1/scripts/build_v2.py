#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Rellena la plantilla UC (10×5,625 in) con el contenido del video de avance 1."""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1F, 0x49, 0x7D)
UCBLUE = RGBColor(0x33, 0x66, 0xCC)
DARK = RGBColor(0x17, 0x37, 0x5D)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF2, 0xF6, 0xFC)
CARD2 = RGBColor(0xF5, 0xF5, 0xF2)
ICE = RGBColor(0xEA, 0xF2, 0xFD)
GOLD = RGBColor(0xFF, 0xE2, 0x00)
ORANGE = RGBColor(0xB2, 0x4E, 0x1C)
AQUA = RGBColor(0x0F, 0x7D, 0x57)
LIGHTBLUE = RGBColor(0xCD, 0xE2, 0xFB)
BORDER = RGBColor(0xD8, 0xDE, 0xE9)

F = 'Calibri'
A = lambda p: f'/home/claude/airppt/assets/{p}'

REFS = {
  'gbd':    'GBD 2021 Risk Factors Collaborators (Brauer M. et al.). Global burden and strength of evidence for 88 risk factors. The Lancet 403 (2024). doi:10.1016/S0140-6736(24)00933-4',
  'who':    'Organización Mundial de la Salud. WHO global air quality guidelines: PM, O3, NO2, SO2 y CO. Ginebra (2021). ISBN 978-92-4-003422-8',
  'villa':  'Villalobos A.M. et al. Wood burning pollution in southern Chile: PM2.5 source apportionment using CMB and molecular markers. Environmental Pollution 225 (2017). doi:10.1016/j.envpol.2017.02.069',
  'barraza':'Barraza F. et al. Temporal evolution of main ambient PM2.5 sources in Santiago, Chile, 1998–2012. Atmos. Chem. Phys. 17 (2017). doi:10.5194/acp-17-10093-2017',
  'vand':   'van Donkelaar A. et al. Monthly global estimates of fine particulate matter and their uncertainty. Environ. Sci. Technol. 55 (2021). doi:10.1021/acs.est.1c05309',
  'larkin': 'Larkin A. et al. Global land use regression model for nitrogen dioxide air pollution. Environ. Sci. Technol. 51 (2017). doi:10.1021/acs.est.7b01148',
  'wei':    'Wei J. et al. Ground-level NO2 surveillance from space across China using interpretable spatiotemporally weighted artificial intelligence. Environ. Sci. Technol. 56 (2022). doi:10.1021/acs.est.2c03834',
  'anen':   'Anenberg S.C. et al. Long-term trends in urban NO2 concentrations and associated paediatric asthma incidence. The Lancet Planetary Health 6 (2022). doi:10.1016/S2542-5196(21)00255-2',
  'turner': 'Turner M.C. et al. Long-term ozone exposure and mortality in a large prospective study. Am. J. Respir. Crit. Care Med. 193 (2016). doi:10.1164/rccm.201508-1633OC',
  'tropomi':'Veefkind J.P. et al. TROPOMI on the ESA Sentinel-5 Precursor. Remote Sensing of Environment 120 (2012). doi:10.1016/j.rse.2011.09.027',
  'omi':    'Levelt P.F. et al. The Ozone Monitoring Instrument. IEEE Trans. Geosci. Remote Sens. 44 (2006). doi:10.1109/TGRS.2006.872333',
  'mopitt': 'Deeter M.N. et al. The MOPITT Version 9 CO product. Atmos. Meas. Tech. 15 (2022). doi:10.5194/amt-15-2325-2022',
  'cams':   'Inness A. et al. The CAMS reanalysis of atmospheric composition. Atmos. Chem. Phys. 19 (2019). doi:10.5194/acp-19-3515-2019',
  'geoscf': 'Keller C.A. et al. Description of the NASA GEOS-CF v1.0. J. Adv. Model. Earth Syst. 13 (2021). doi:10.1029/2020MS002413',
  'merra2': 'Gelaro R. et al. MERRA-2. J. Climate 30 (2017). doi:10.1175/JCLI-D-16-0758.1',
  'era5':   'Hersbach H. et al. The ERA5 global reanalysis. Q. J. R. Meteorol. Soc. 146 (2020). doi:10.1002/qj.3803',
  'maiac':  'Lyapustin A. et al. MODIS Collection 6 MAIAC algorithm. Atmos. Meas. Tech. 11 (2018). doi:10.5194/amt-11-5741-2018',
  'norma':  'Ministerio del Medio Ambiente, Chile. D.S. N.º 12/2011, norma primaria de calidad ambiental para MP2,5.',
}
SUP = ['¹','²','³','⁴','⁵','⁶','⁷','⁸','⁹']

prs = Presentation('base.pptx')
S = prs.slides


def no_line(shape):
    shape.line.fill.background()


def no_shadow(shape):
    el = shape._element.spPr
    existing = el.find(qn('a:effectLst'))
    if existing is None:
        from lxml import etree
        el.append(etree.SubElement(el, qn('a:effectLst')))


def box(slide, x, y, w, h, fill, line=None, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.adjustments[0] = radius
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    no_shadow(sp)
    sp.text_frame.paragraphs[0].text = ''
    return sp


def oval(slide, x, y, d, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    no_line(sp); no_shadow(sp)
    return sp


def icono(slide, name, x, y, d, circ):
    oval(slide, x, y, d, circ)
    pad = d * 0.24
    slide.shapes.add_picture(A(f'icons/{name}_w.png'), Inches(x + pad), Inches(y + pad), Inches(d - 2 * pad), Inches(d - 2 * pad))


def text(slide, x, y, w, h, runs, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, leading=1.0, space_after=0, wrap=True):
    """runs: str o lista de (texto, dict-overrides) o lista de párrafos (lista de runs)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if leading != 1.0:
            p.line_spacing = leading
        if space_after:
            p.space_after = Pt(space_after)
        for t, ov in para:
            r = p.add_run(); r.text = t
            r.font.name = F
            r.font.size = Pt(ov.get('size', size))
            r.font.bold = ov.get('bold', bold)
            r.font.italic = ov.get('italic', False)
            r.font.color.rgb = ov.get('color', color)
    return tb


def bullets(slide, x, y, w, h, items, size=9, color=INK2, gap=4, leading=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if leading != 1.0:
            p.line_spacing = leading
        r = p.add_run(); r.text = '• ' + it
        r.font.name = F; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def footnotes(slide, keys, y, w=7.05, x=0.5, size=6.2, two_cols=False):
    def block(ks, sup0, bx, bw):
        paras = [[(f'{SUP[sup0 + i]} {REFS[k]}', {})] for i, k in enumerate(ks)]
        text(slide, bx, y, bw, 5.1 - y, paras, size=size, color=MUTED, leading=1.0, space_after=1)
    if two_cols:
        mid = (len(keys) + 1) // 2
        block(keys[:mid], 0, x, 4.45)
        block(keys[mid:], mid, x + 4.6, 4.4)
    else:
        block(keys, 0, x, w)


def set_title(slide, texto, size=20):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 or 'Título' in ph.name or 'Title' in ph.name:
            tf = ph.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = texto
            r.font.name = F; r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = WHITE
            return ph


def drop_content_placeholders(slide):
    for ph in list(slide.placeholders):
        if 'contenido' in ph.name.lower():
            ph._element.getparent().remove(ph._element)


# =====================================================================
# S2 (índice 0): portada de título
# =====================================================================
s = S[0]
for ph in s.placeholders:
    if 'Título' in ph.name:
        tf = ph.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = 'Estimación multi-contaminante de la calidad del aire en Chile'
        r.font.name = F; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = DARK
    elif 'Subtítulo' in ph.name:
        ph.left = Inches(0.67); ph.width = Inches(8.5); ph.top = Inches(3.02); ph.height = Inches(1.2)
        tf = ph.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ('Superficies de exposición comuna × hora para PM₂.₅ · PM₁₀ · NO₂ · O₃ · SO₂ · CO, '
                  'para todo el territorio nacional y por macrozona, validadas contra la red SINCA')
        r.font.name = F; r.font.size = Pt(13); r.font.color.rgb = INK2
text(s, 0.67, 1.30, 8.5, 0.3, [( 'ACTIVIDAD FINAL DE GRADUACIÓN 1  ·  MDS3050  ·  VIDEO DE AVANCE N.º 1: DEFINICIÓN DEL PROBLEMA', {})],
     size=10.5, color=UCBLUE, bold=True)
equipo = ['José Jesús Romero Fuenmayor', 'Roberto Ignacio Ávila Escobar', 'Amaru Simón Agüero Jiménez']
for i, n in enumerate(equipo):
    x = 0.67 + i * 2.92
    box(s, x, 4.32, 2.78, 0.44, ICE, radius=0.18)
    text(s, x + 0.12, 4.32, 2.6, 0.44, n, size=10, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.67, 4.88, 8.5, 0.26, 'Equipo 2 (Team2)  ·  agosto de 2026', size=9.5, color=MUTED)

# =====================================================================
# S3 (índice 1): ¿por qué importa?
# =====================================================================
s = S[1]
set_title(s, '¿Por qué importa? El aire que se respira mata y casi no se mide')
drop_content_placeholders(s)
stats = [
    ('≈ 8 millones', 'de muertes al año se atribuyen a la contaminación del aire en el mundo (GBD 2021)¹', 'health', NAVY),
    ('4 de 6', 'contaminantes normados endurecieron su guía OMS 2021²; el NO₂ anual pasó de 40 a 10 µg/m³', 'warning', UCBLUE),
    ('4 ×', 'la norma anual chilena de PM₂.₅ cuadruplica la guía OMS vigente² ³ (20 frente a 5 µg/m³)', 'map', AQUA),
]
for i, (big, small, ic, c) in enumerate(stats):
    x = 0.5 + i * 3.07
    box(s, x, 1.22, 2.93, 1.58, CARD)
    icono(s, ic, x + 0.16, 1.36, 0.44, c)
    text(s, x + 0.70, 1.30, 2.15, 0.55, big, size=20, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.18, 1.92, 2.6, 0.82, small, size=8.3, color=INK2, leading=0.95)
text(s, 0.5, 2.94, 9.0, 0.26, 'En Chile el problema tiene dos caras bien documentadas', size=11.5, bold=True, color=NAVY)
caras = [
    ('Centro-sur: leña residencial', 'La quema de leña domina el PM₂.₅ invernal en ciudades como Temuco; el problema es de partículas y de los gases que las acompañan.⁴', ORANGE),
    ('Santiago y zonas industriales', 'Fuentes vehiculares e industriales sostienen episodios de PM₂.₅, NO₂ y SO₂; su composición cambió durante 15 años de gestión.⁵', UCBLUE),
]
for i, (t, d, c) in enumerate(caras):
    x = 0.5 + i * 4.62
    box(s, x, 3.26, 4.38, 1.06, WHITE, line=BORDER)
    text(s, x + 0.16, 3.36, 4.06, 0.24, t, size=9.5, bold=True, color=c)
    text(s, x + 0.16, 3.62, 4.06, 0.66, d, size=8.2, color=INK2, leading=0.95)
footnotes(s, ['gbd', 'who', 'norma', 'villa', 'barraza'], y=4.46)

# =====================================================================
# S4 (índice 2): el problema
# =====================================================================
s = S[2]
set_title(s, 'El problema: los gases casi no se vigilan fuera de las grandes ciudades', size=18)
drop_content_placeholders(s)
box(s, 0.5, 1.22, 5.55, 2.34, ICE, line=UCBLUE)
text(s, 0.72, 1.36, 5.1, 0.24, 'FORMULACIÓN', size=9, bold=True, color=NAVY)
text(s, 0.72, 1.64, 5.14, 2.1,
     'La vigilancia de calidad del aire en Chile descansa en ~109 estaciones SINCA concentradas en '
     'centros urbanos; para los gases (NO₂, O₃, SO₂, CO) la cobertura es aún menor que para el material '
     'particulado. No existe una superficie pública, continua y validada de estos contaminantes para '
     'todo el territorio. Proponemos estimarla a resolución comuna × hora / día combinando satélites, '
     'reanálisis y predictores locales, con la red SINCA como única verdad-terreno, y reportarla a '
     'escala nacional y por macrozona.', size=9.6, color=INK, leading=1.02)
lat = [
    ('Unidad de análisis', 'comuna × hora (y día), 345 comunas; agregación por macrozona', 'map'),
    ('Horizonte temporal', 'según fuente: MOPITT 2000+ · OMI 2004+ · CAMS 2003+ · TROPOMI / GEOS-CF 2018+', 'calendar'),
    ('Alcance acotado', 'estimar bien y validar contra SINCA; sin validación epidemiológica (recomendación docente)', 'target'),
]
for i, (t, d, ic) in enumerate(lat):
    y = 1.22 + i * 0.92
    box(s, 6.22, y, 3.28, 0.80, CARD)
    icono(s, ic, 6.34, y + 0.16, 0.46, NAVY)
    text(s, 6.94, y + 0.08, 2.48, 0.22, t, size=9, bold=True, color=NAVY)
    text(s, 6.94, y + 0.30, 2.48, 0.48, d, size=7.4, color=INK2, leading=0.92)
box(s, 0.5, 3.96, 7.08, 0.98, WHITE, line=BORDER)
icono(s, 'download', 0.68, 4.19, 0.52, AQUA)
text(s, 1.38, 4.06, 6.02, 0.82,
     [(( 'Punto de partida operativo. '), {'bold': True, 'color': INK}),
      (('El repositorio del proyecto ya integra descargadores verificados para 10 productos '
        'satelitales / de reanálisis y la red SINCA, con recorte automático a Chile, estructura '
        'reproducible y bibliografía en formato APA.'), {'color': INK2})],
     size=9.3, leading=1.0, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# S5 (índice 3): evidencia
# =====================================================================
s = S[3]
set_title(s, 'La evidencia dice que se puede: precedentes internacionales')
drop_content_placeholders(s)
cards = [
    ('PM₂.₅ global mensual', 'van Donkelaar 2021', 'Estimaciones globales de PM₂.₅ con incertidumbre, combinando satélite + modelo + estaciones (base del producto ACAG que usamos).¹'),
    ('NO₂ global (LUR)', 'Larkin 2017', 'Modelo global de regresión de uso de suelo para NO₂ anual: el precedente de escalar gases a territorio completo.²'),
    ('NO₂ diario con IA', 'Wei 2022', 'NO₂ de superficie sin huecos para toda China con IA espaciotemporal interpretable; el estándar técnico al que apuntamos.³'),
    ('Gases → decisiones', 'Anenberg 2022 · Turner 2016', 'Con esas superficies se estiman impactos: tendencias urbanas de NO₂ y asma pediátrica⁴; el O₃ de largo plazo se asocia a mortalidad.⁵'),
]
for i, (t, tag, d) in enumerate(cards):
    col, row = i % 2, i // 2
    x, y = 0.5 + col * 4.62, 1.22 + row * 1.38
    box(s, x, y, 4.38, 1.24, CARD)
    text(s, x + 0.16, y + 0.10, 2.6, 0.24, t, size=10.5, bold=True, color=NAVY)
    box(s, x + 2.62, y + 0.10, 1.62, 0.24, ICE, radius=0.3)
    text(s, x + 2.62, y + 0.10, 1.62, 0.24, tag, size=6.6, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.16, y + 0.40, 4.06, 0.8, d, size=8.2, color=INK2, leading=0.96)
text(s, 0.5, 4.06, 9.0, 0.5,
     [(('Brecha que abordamos: '), {'bold': True, 'color': NAVY}),
      (('para Chile no existe un producto multi-gas continuo, público y validado por estación; los '
        'precedentes existen para PM₂.₅ global y para gases en el hemisferio norte.'), {'color': INK2})],
     size=9.3, leading=1.0)
footnotes(s, ['vand', 'larkin', 'wei', 'anen', 'turner'], y=4.56)

# =====================================================================
# S6 (índice 4): objetivos
# =====================================================================
s = S[4]
set_title(s, 'Objetivo general y objetivos específicos')
drop_content_placeholders(s)
box(s, 0.5, 1.20, 9.0, 1.02, DARK)
icono(s, 'target', 0.70, 1.42, 0.58, UCBLUE)
text(s, 1.48, 1.30, 7.85, 0.84,
     [(('OBJETIVO GENERAL: '), {'bold': True, 'color': GOLD}),
      (('Generar y validar superficies de exposición de PM₂.₅, PM₁₀, NO₂, O₃, SO₂ y CO a resolución '
        'comuna × hora / día para todo Chile, agregadas por macrozona, mediante un pipeline reproducible '
        'que integra satélites, reanálisis y predictores locales, evaluado contra la red SINCA.'), {'color': WHITE})],
     size=9.6, leading=1.0, anchor=MSO_ANCHOR.MIDDLE)
oes = [
    ('OE1', 'Consolidar el pipeline de descarga y procesamiento multi-fuente (10 productos + SINCA), con recorte a Chile y agregación comunal.', 'download'),
    ('OE2', 'Construir paneles de modelamiento por contaminante uniendo predictores satelitales, meteorológicos y estáticos con SINCA.', 'storage'),
    ('OE3', 'Entrenar y comparar motores de estimación (GWR, gradient boosting, kriging) con validación cruzada LOSO.', 'science'),
    ('OE4', 'Producir las superficies comuna × hora / día y reportar métricas (R², RMSE, sesgo) por contaminante, estación y macrozona.', 'map'),
    ('OE5', 'Publicar el proyecto como repositorio abierto y reproducible (scripts, credenciales de ejemplo, documentación).', 'checklist'),
]
for i, (n, t, ic) in enumerate(oes):
    col, row = i % 2, i // 2
    x, y = 0.5 + col * 4.62, 2.38 + row * 0.98
    box(s, x, y, 4.38, 0.88, CARD)
    icono(s, ic, x + 0.12, y + 0.20, 0.46, UCBLUE)
    text(s, x + 0.70, y + 0.07, 0.7, 0.22, n, size=9.5, bold=True, color=NAVY)
    text(s, x + 0.70, y + 0.28, 3.58, 0.56, t, size=7.6, color=INK2, leading=0.92)
text(s, 5.12, 4.40, 4.38, 0.88,
     [(('Hipótesis de trabajo: '), {'bold': True, 'italic': True, 'color': NAVY}),
      (('los predictores satelitales y de reanálisis permiten estimar los gases de superficie con '
        'desempeño comparable al reportado internacionalmente, incluso en macrozonas con pocas '
        'estaciones.'), {'italic': True, 'color': INK2})],
     size=8.4, leading=1.0)

# =====================================================================
# S7 (índice 5): flujograma
# =====================================================================
s = S[5]
set_title(s, 'Flujograma: de dónde salen los datos y cómo se transforman')
drop_content_placeholders(s)
from PIL import Image
iw, ih = Image.open(A('flujograma_wide.png')).size
ar = iw / ih
y0 = 1.18
w = min(8.70, (5.02 - y0) * ar)
h = w / ar
s.shapes.add_picture(A('flujograma_wide.png'), Inches((10 - w) / 2), Inches(y0), Inches(w), Inches(h))

# =====================================================================
# S8 (índice 6): tabla de fuentes
# =====================================================================
s = S[6]
set_title(s, 'Fuentes de datos por contaminante', size=19)
drop_content_placeholders(s)
rows = [
    ['Contaminante', 'Satélite (columna / L2–L3)', 'Reanálisis / superficie', 'Cobertura'],
    ['PM₂.₅', 'MAIAC AOD¹ · MODIS AOD 3K (predictores)', 'ACAG² · MERRA-2³ · CAMS⁴ · GEOS-CF⁵', '2000+'],
    ['PM₁₀', 'MODIS AOD (predictor)', 'CAMS⁴ · MERRA-2 (polvo)³', '2003+'],
    ['NO₂', 'TROPOMI⁶ · OMI⁷', 'CAMS⁴ · GEOS-CF⁵', '2004+'],
    ['O₃', 'TROPOMI⁶ · OMI⁷', 'MERRA-2³ · CAMS⁴ · GEOS-CF⁵', '2004+'],
    ['SO₂', 'TROPOMI⁶ · OMI⁷', 'CAMS⁴ · GEOS-CF⁵', '2004+'],
    ['CO', 'TROPOMI⁶ · MOPITT⁸', 'CAMS⁴ · GEOS-CF⁵', '2000+'],
]
tw = [1.25, 3.55, 3.15, 1.05]
tbl_shape = s.shapes.add_table(7, 4, Inches(0.5), Inches(1.22), Inches(9.0), Inches(2.32))
tbl = tbl_shape.table
tbl.first_row = True; tbl.horz_banding = True
for j, wcol in enumerate(tw):
    tbl.columns[j].width = Inches(wcol)
for i in range(7):
    tbl.rows[i].height = Inches(0.33)
    for j in range(4):
        c = tbl.cell(i, j)
        c.margin_left = Inches(0.06); c.margin_right = Inches(0.04)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = rows[i][j]
        r.font.name = F
        if i == 0:
            r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = WHITE
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
        else:
            r.font.size = Pt(8.5); r.font.color.rgb = INK
            r.font.bold = (j == 0)
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else WHITE
text(s, 0.5, 3.72, 9.0, 0.5,
     [(('Predictores comunes a los 6: '), {'bold': True, 'color': NAVY}),
      (('meteorología ERA5 / ERA5-Land⁹ y MERRA-2³, uso de suelo ESA CCI, red vial OSM, altitud '
        'NASADEM, luces nocturnas VIIRS y censo (población, leña).'), {'color': INK2})],
     size=8.8, leading=1.0)
footnotes(s, ['maiac', 'vand', 'merra2', 'cams', 'geoscf', 'tropomi', 'omi', 'mopitt', 'era5'],
          y=4.28, two_cols=True)

# =====================================================================
# S9 (índice 7): semana
# =====================================================================
s = S[7]
set_title(s, 'Semanas 1 y 2: qué hicimos, qué costó y qué viene')
drop_content_placeholders(s)
colsx = [
    ('Objetivos de la semana', 'checklist', NAVY, [
        'Plantear la problemática real como proyecto de ciencia de datos',
        'Delimitar objetivo general y específicos',
        'Dejar operativa la obtención de datos multi-fuente']),
    ('Tareas realizadas', 'download', UCBLUE, [
        'Repositorio público con estructura reproducible y bibliografía APA',
        'Descargadores con IDs verificados en CMR / ADS / S3',
        'Recorte a Chile (~5 % del volumen) y agregación comunal',
        'Config SINCA: 109 estaciones, 6 contaminantes']),
    ('Desafíos', 'warning', ORANGE, [
        'Volumen: TROPOMI sin recorte ≈ 13 TB; MERRA-2 ≈ 1 TB',
        'Credenciales y cuotas (Earthdata, ADS) y reintentos',
        'Cobertura temporal heterogénea (2000+ / 2004+ / 2018+)']),
    ('Próxima semana', 'calendar', AQUA, [
        'Revisión de literatura por contaminante (video N.º 2)',
        'Métricas objetivo y líneas base por macrozona',
        'Preparar la defensa de tema (semana 4)']),
]
for i, (t, ic, c, items) in enumerate(colsx):
    x = 0.5 + i * 2.32
    box(s, x, 1.22, 2.20, 2.82, CARD)
    icono(s, ic, x + 0.14, 1.36, 0.42, c)
    text(s, x + 0.64, 1.34, 1.5, 0.5, t, size=9.5, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE, leading=0.9)
    bullets(s, x + 0.16, 1.98, 1.9, 2.0, items, size=7.6, gap=4, leading=0.95)

# =====================================================================
# S10 (índice 8): cierre con QR (fondo azul de la plantilla)
# =====================================================================
s = S[8]
text(s, 0.55, 1.02, 5.9, 0.75, 'Todo el proyecto vive aquí,\nabierto y reproducible', size=19, bold=True, color=WHITE, leading=1.0)
bullets(s, 0.55, 1.95, 5.9, 1.7, [
    'scripts_pipeline/: descarga y procesamiento por contaminante',
    'scripts_superficie/: modelos, validación LOSO y figuras',
    'docs/: sitio con bibliografía APA consolidada',
    'Datos pesados regenerables con credenciales propias (Earthdata / ADS)',
], size=9.5, color=LIGHTBLUE, gap=5)
box(s, 0.55, 3.78, 5.9, 0.56, RGBColor(0x11, 0x3E, 0x8F))
s.shapes.add_picture(A('icons/github_w.png'), Inches(0.72), Inches(3.92), Inches(0.28), Inches(0.28))
text(s, 1.10, 3.78, 5.3, 0.56, 'github.com/AmaruSimonAgueroJimenez/Air-Pollution', size=11.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.95, 1.02, 2.55, 3.32, WHITE, radius=0.06)
s.shapes.add_picture(A('qr_repo.png'), Inches(7.17), Inches(1.24), Inches(2.11), Inches(2.11))
text(s, 7.05, 3.42, 2.35, 0.8, 'Escanea para ver el repositorio, el pipeline y la documentación',
     size=8.5, color=DARK, align=PP_ALIGN.CENTER, leading=0.95)
text(s, 0.55, 4.80, 8.9, 0.3, 'Equipo 2: José J. Romero · Roberto I. Ávila · Amaru S. Agüero          ¡Gracias!',
     size=10, bold=False, color=LIGHTBLUE)

prs.save('AFG1_Video1_Presentacion_UC.pptx')
print('ok')
