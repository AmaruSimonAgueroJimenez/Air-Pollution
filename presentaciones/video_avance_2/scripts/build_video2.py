#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Video de avance N.º 2 (revisión de literatura) sobre la plantilla UC."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1F, 0x49, 0x7D)
UCBLUE = RGBColor(0x33, 0x66, 0xCC)
DARK = RGBColor(0x17, 0x37, 0x5D)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF2, 0xF6, 0xFC)
ICE = RGBColor(0xEA, 0xF2, 0xFD)
GOLD = RGBColor(0xFF, 0xE2, 0x00)
ORANGE = RGBColor(0xB2, 0x4E, 0x1C)
AQUA = RGBColor(0x0F, 0x7D, 0x57)
LIGHTBLUE = RGBColor(0xCD, 0xE2, 0xFB)
BORDER = RGBColor(0xD8, 0xDE, 0xE9)
CREAM = RGBColor(0xFD, 0xF2, 0xEC)

F = 'Calibri'
A = lambda p: f'/home/claude/airppt/assets/{p}'

REFS = {
  'vand':   'van Donkelaar A. et al. Monthly global estimates of fine particulate matter and their uncertainty. Environ. Sci. Technol. 55 (2021). doi:10.1021/acs.est.1c05309',
  'vd16':   'van Donkelaar A. et al. Global estimates of fine particulate matter using a combined geophysical-statistical method with information from satellites, models, and monitors. Environ. Sci. Technol. 50 (2016). doi:10.1021/acs.est.5b05833',
  'di19':   'Di Q. et al. An ensemble-based model of PM2.5 concentration across the contiguous United States with high spatiotemporal resolution. Environment International (2019). doi:10.1016/j.envint.2019.104909',
  'villa':  'Villalobos A.M. et al. Wood burning pollution in southern Chile: PM2.5 source apportionment using CMB and molecular markers. Environmental Pollution 225 (2017). doi:10.1016/j.envpol.2017.02.069',
  'barraza':'Barraza F. et al. Temporal evolution of main ambient PM2.5 sources in Santiago, Chile, 1998–2012. Atmos. Chem. Phys. 17 (2017). doi:10.5194/acp-17-10093-2017',
  'larkin': 'Larkin A. et al. Global land use regression model for nitrogen dioxide air pollution. Environ. Sci. Technol. 51 (2017). doi:10.1021/acs.est.7b01148',
  'wei':    'Wei J. et al. Ground-level NO2 surveillance from space across China using interpretable spatiotemporally weighted artificial intelligence. Environ. Sci. Technol. 56 (2022). doi:10.1021/acs.est.2c03834',
  'anen':   'Anenberg S.C. et al. Long-term trends in urban NO2 concentrations and associated paediatric asthma incidence. The Lancet Planetary Health 6 (2022). doi:10.1016/S2542-5196(21)00255-2',
  'turner': 'Turner M.C. et al. Long-term ozone exposure and mortality in a large prospective study. Am. J. Respir. Crit. Care Med. 193 (2016). doi:10.1164/rccm.201508-1633OC',
  'wei23':  'Wei J. et al. Ground-level gaseous pollutants (NO2, SO2, and CO) in China: daily seamless mapping and spatiotemporal variations. Atmos. Chem. Phys. 23, 1511–1532 (2023). doi:10.5194/acp-23-1511-2023',
  'theys':  'Theys N. et al. Sulfur dioxide retrievals from TROPOMI onboard Sentinel-5 Precursor: algorithm theoretical basis. Atmos. Meas. Tech. 10, 119–153 (2017). doi:10.5194/amt-10-119-2017',
  'mopitt': 'Deeter M.N. et al. The MOPITT Version 9 CO product. Atmos. Meas. Tech. 15 (2022). doi:10.5194/amt-15-2325-2022',
  'hoek':   'Hoek G. et al. A review of land-use regression models to assess spatial variation of outdoor air pollution. Atmospheric Environment 42 (2008). doi:10.1016/j.atmosenv.2008.05.057',
  'cams':   'Inness A. et al. The CAMS reanalysis of atmospheric composition. Atmos. Chem. Phys. 19 (2019). doi:10.5194/acp-19-3515-2019',
  'geoscf': 'Keller C.A. et al. Description of the NASA GEOS-CF v1.0. J. Adv. Model. Earth Syst. 13 (2021). doi:10.1029/2020MS002413',
}
SUP = ['¹','²','³','⁴','⁵','⁶','⁷','⁸','⁹']

prs = Presentation('base.pptx')
S = prs.slides


def no_line(shape):
    shape.line.fill.background()


def no_shadow(shape):
    el = shape._element.spPr
    if el.find(qn('a:effectLst')) is None:
        from lxml import etree
        el.append(etree.SubElement(el, qn('a:effectLst')))


def box(slide, x, y, w, h, fill, line=None, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.adjustments[0] = radius
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    no_shadow(sp)
    sp.text_frame.paragraphs[0].text = ''
    return sp


def oval(slide, x, y, d, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    no_line(sp); no_shadow(sp)
    return sp


def icono(slide, name, x, y, d, circ):
    oval(slide, x, y, d, circ)
    pad = d * 0.24
    slide.shapes.add_picture(A(f'icons/{name}_w.png'), Inches(x + pad), Inches(y + pad), Inches(d - 2 * pad), Inches(d - 2 * pad))


def text(slide, x, y, w, h, runs, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, leading=1.0, space_after=0, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if leading != 1.0:
            p.line_spacing = leading
        if space_after:
            p.space_after = Pt(space_after)
        for t, ov in para:
            r = p.add_run(); r.text = t
            r.font.name = F
            r.font.size = Pt(ov.get('size', size))
            r.font.bold = ov.get('bold', bold)
            r.font.italic = ov.get('italic', False)
            r.font.color.rgb = ov.get('color', color)
    return tb


def bullets(slide, x, y, w, h, items, size=9, color=INK2, gap=4, leading=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if leading != 1.0:
            p.line_spacing = leading
        r = p.add_run(); r.text = '• ' + it
        r.font.name = F; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def footnotes(slide, keys, y, w=7.05, x=0.5, size=6.2, two_cols=False):
    def block(ks, sup0, bx, bw):
        paras = [[(f'{SUP[sup0 + i]} {REFS[k]}', {})] for i, k in enumerate(ks)]
        text(slide, bx, y, bw, 5.35 - y, paras, size=size, color=MUTED, leading=1.0, space_after=1)
    if two_cols:
        mid = (len(keys) + 1) // 2
        block(keys[:mid], 0, x, 4.45)
        block(keys[mid:], mid, x + 4.6, 4.4)
    else:
        block(keys, 0, x, w)


def set_title(slide, texto, size=20):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 or 'Título' in ph.name or 'Title' in ph.name:
            tf = ph.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = texto
            r.font.name = F; r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = WHITE
            return ph


def drop_content_placeholders(slide):
    for ph in list(slide.placeholders):
        if 'contenido' in ph.name.lower():
            ph._element.getparent().remove(ph._element)


# =====================================================================
# S0: portada
# =====================================================================
s = S[0]
for ph in s.placeholders:
    if 'Título' in ph.name:
        tf = ph.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = 'Revisión de literatura: estimar la calidad del aire desde el espacio'
        r.font.name = F; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = DARK
    elif 'Subtítulo' in ph.name:
        ph.left = Inches(0.67); ph.width = Inches(8.5); ph.top = Inches(3.02); ph.height = Inches(1.2)
        tf = ph.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ('Qué se sabe, con qué métodos se ha hecho y qué brecha abordamos para estimar '
                  'PM₂.₅ · PM₁₀ · NO₂ · O₃ · SO₂ · CO en todo Chile, validado contra la red SINCA')
        r.font.name = F; r.font.size = Pt(13); r.font.color.rgb = INK2
text(s, 0.67, 1.30, 8.5, 0.3, [( 'ACTIVIDAD FINAL DE GRADUACIÓN 1  ·  MDS3050  ·  VIDEO DE AVANCE N.º 2: REVISIÓN DE LITERATURA', {})],
     size=10.5, color=UCBLUE, bold=True)
equipo = ['José Jesús Romero Fuenmayor', 'Roberto Ignacio Ávila Escobar',
          'Amaru Simón Agüero Jiménez', 'Esteban Adolfo González Rodríguez']
for i, n in enumerate(equipo):
    x = 0.67 + i * 2.20
    box(s, x, 4.32, 2.06, 0.44, ICE, radius=0.18)
    text(s, x + 0.08, 4.32, 1.90, 0.44, n, size=8.2, bold=True, color=DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.67, 4.88, 8.5, 0.26, 'Equipo 2 (Team2)  ·  agosto de 2026', size=9.5, color=MUTED)

# =====================================================================
# S1: cómo hicimos la revisión
# =====================================================================
s = S[1]
set_title(s, 'Cómo hicimos la revisión: método antes que lista de papers')
drop_content_placeholders(s)
pasos = [
    ('Búsqueda por palabras clave', 'Google Scholar, Scopus y Web of Science; términos en inglés por contaminante: satellite, surface estimation, LUR, machine learning', 'explore', NAVY),
    ('Criterios de inclusión', 'trabajos 2000+; estimación de superficies con validación contra estaciones; guías OMS y normas como literatura oficial', 'checklist', UCBLUE),
    ('Citation chaining', 'desde los trabajos ancla, hacia atrás (sus referencias) y hacia adelante (quién los cita hoy)', 'science', AQUA),
    ('Síntesis y gestión', 'matriz contaminante × método; bibliografía en formato APA consolidada en el repositorio', 'storage', ORANGE),
]
for i, (t, d, ic, c) in enumerate(pasos):
    x = 0.5 + i * 2.32
    box(s, x, 1.22, 2.20, 2.44, CARD)
    icono(s, ic, x + 0.14, 1.36, 0.42, c)
    text(s, x + 0.64, 1.34, 1.5, 0.5, t, size=9.5, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE, leading=0.9)
    text(s, x + 0.16, 2.02, 1.9, 1.55, d, size=7.9, color=INK2, leading=0.98)
box(s, 0.5, 3.86, 9.0, 0.62, ICE)
text(s, 0.68, 3.94, 8.65, 0.5,
     [(('Estructura elegida: '), {'bold': True, 'color': NAVY}),
      (('por tema (cada contaminante) y por metodología, dos de las formas recomendadas en la sesión '
        'sincrónica del curso. La revisión busca qué se sabe, cómo se llegó y qué no se sabe.'), {'color': INK2})],
     size=9.2, leading=1.0)
text(s, 0.5, 4.62, 9.0, 0.3, 'Todas las referencias citadas en esta presentación tienen DOI o ISBN verificado y están en la bibliografía del repositorio.',
     size=8.2, color=MUTED)

# =====================================================================
# S2: qué se sabe PM2.5
# =====================================================================
s = S[2]
set_title(s, 'Qué se sabe (1): PM₂.₅ es el contaminante mejor resuelto', size=19)
drop_content_placeholders(s)
cards = [
    ('Global, mensual y con incertidumbre', 'van Donkelaar 2021 · ACAG', 'Estimaciones globales de PM₂.₅ combinando satélite, modelo químico y estaciones; es la base del producto de superficie que usamos.¹'),
    ('El salto: híbrido geofísico + estadístico', 'van Donkelaar 2016', 'AOD satelital + GEOS-Chem, corregido con regresión geográficamente ponderada contra estaciones: R² 0.81 en validación cruzada global.²'),
    ('Ensambles de aprendizaje automático', 'Di 2019', 'EE. UU. a 1 km y resolución diaria con un ensamble de redes neuronales, boosting y bosques aleatorios: el estándar de ingeniería actual.³'),
    ('Chile: composición y fuentes locales', 'Villalobos 2017 · Barraza 2017', 'La leña domina el PM₂.₅ invernal del centro-sur⁴ y las fuentes de Santiago evolucionaron durante 15 años⁵: contexto para interpretar nuestras superficies.'),
]
for i, (t, tag, d) in enumerate(cards):
    col, row = i % 2, i // 2
    x, y = 0.5 + col * 4.62, 1.22 + row * 1.38
    box(s, x, y, 4.38, 1.24, CARD)
    text(s, x + 0.16, y + 0.10, 2.55, 0.4, t, size=9.8, bold=True, color=NAVY, leading=0.9)
    box(s, x + 2.72, y + 0.10, 1.52, 0.24, ICE, radius=0.3)
    text(s, x + 2.72, y + 0.10, 1.52, 0.24, tag, size=6.2, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.16, y + 0.48, 4.06, 0.72, d, size=8.0, color=INK2, leading=0.95)
text(s, 0.5, 4.06, 9.0, 0.3,
     [(('La receta madura: '), {'bold': True, 'color': NAVY}),
      (('AOD satelital + modelo químico + predictores locales + calibración contra estaciones.'), {'color': INK2})],
     size=9.3)
footnotes(s, ['vand', 'vd16', 'di19', 'villa', 'barraza'], y=4.42)

# =====================================================================
# S3: qué se sabe gases
# =====================================================================
s = S[3]
set_title(s, 'Qué se sabe (2): los gases van una década más atrás', size=19)
drop_content_placeholders(s)
box(s, 0.5, 1.22, 4.42, 2.60, CARD)
text(s, 0.68, 1.34, 4.1, 0.26, 'NO₂ y O₃: el camino ya está trazado', size=10.5, bold=True, color=UCBLUE)
bullets(s, 0.68, 1.68, 4.06, 2.0, [
    'LUR global de NO₂ anual: el precedente de escalar un gas a todo el planeta¹',
    'NO₂ diario sin huecos para China con IA espaciotemporal interpretable²',
    'Las superficies alimentan decisiones: asma pediátrica por NO₂³ y mortalidad por O₃ de largo plazo⁴',
], size=8.4, gap=5, leading=0.98)
box(s, 5.08, 1.22, 4.42, 2.60, CREAM, line=ORANGE)
text(s, 5.26, 1.34, 4.1, 0.26, 'SO₂ y CO: los menos estudiados', size=10.5, bold=True, color=ORANGE)
bullets(s, 5.26, 1.68, 4.06, 2.0, [
    'Primer mapeo diario y sin huecos de NO₂, SO₂ y CO para China recién en 2023⁵',
    'La base instrumental existe: algoritmo SO₂ de TROPOMI⁶ y 20+ años de CO desde MOPITT⁷',
    'Fuera de estación el desempeño cae a R² 0.6 a 0.7⁵: medir esa caída importa',
], size=8.4, gap=5, leading=0.98)
box(s, 0.5, 4.00, 9.0, 0.44, ICE)
text(s, 0.68, 4.06, 8.65, 0.32,
     [(('Patrón geográfico: '), {'bold': True, 'color': NAVY}),
      (('casi toda esta literatura se concentra en China, EE. UU. y Europa; Sudamérica no tiene un producto propio equivalente.'), {'color': INK2})],
     size=9.2)
footnotes(s, ['larkin', 'wei', 'anen', 'turner', 'wei23', 'theys', 'mopitt'], y=4.50, size=5.9, two_cols=True)

# =====================================================================
# S4: metodologías
# =====================================================================
s = S[4]
set_title(s, 'Cómo se ha hecho: cuatro familias de métodos', size=19)
drop_content_placeholders(s)
rows = [
    ['Familia', 'Idea central', 'Ejemplos verificados'],
    ['Regresión de uso de suelo (LUR)', 'predictores de suelo, vías y población explican el gradiente espacial', 'Hoek 2008 (revisión canónica)¹ · Larkin 2017²'],
    ['Geoestadística y GWR', 'coeficientes que varían en el espacio; kriging de residuos', 'van Donkelaar 2016³'],
    ['Aprendizaje automático', 'boosting, bosques y ensambles espaciotemporales', 'Di 2019⁴ · Wei 2022⁵ · Wei 2023⁶'],
    ['Híbridos físico-estadísticos', 'modelo químico + satélite + corrección con estaciones', 'van Donkelaar 2016³ y 2021⁷ · CAMS⁸ / GEOS-CF⁹ como insumo'],
]
tw = [2.30, 3.45, 3.25]
tbl_shape = s.shapes.add_table(5, 3, Inches(0.5), Inches(1.22), Inches(9.0), Inches(2.10))
tbl = tbl_shape.table
tbl.first_row = True; tbl.horz_banding = True
for j, wcol in enumerate(tw):
    tbl.columns[j].width = Inches(wcol)
for i in range(5):
    tbl.rows[i].height = Inches(0.40 if i else 0.30)
    for j in range(3):
        c = tbl.cell(i, j)
        c.margin_left = Inches(0.07); c.margin_right = Inches(0.05)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = rows[i][j]
        r.font.name = F
        if i == 0:
            r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
        else:
            r.font.size = Pt(8.6); r.font.color.rgb = INK
            r.font.bold = (j == 0)
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else WHITE
box(s, 0.5, 3.56, 9.0, 0.50, ICE)
text(s, 0.68, 3.63, 8.65, 0.38,
     [(('Validación estándar del área: '), {'bold': True, 'color': NAVY}),
      (('cross-validation espacial dejando estaciones fuera. Nuestra elección de LOSO y de los motores GWR, '
        'boosting y kriging sale directamente de estas cuatro familias.'), {'color': INK2})],
     size=9.0, leading=0.98)
footnotes(s, ['hoek', 'larkin', 'vd16', 'di19', 'wei', 'wei23', 'vand', 'cams', 'geoscf'], y=4.22, two_cols=True)

# =====================================================================
# S5: brechas y posicionamiento
# =====================================================================
s = S[5]
set_title(s, 'Qué no se sabe: brechas y nuestro posicionamiento', size=19)
drop_content_placeholders(s)
brechas = [
    ('Cobertura geográfica', 'Hay productos multi-gas continuos para China, EE. UU. y Europa; para Chile y Sudamérica no existe un equivalente público y validado.', 'map'),
    ('SO₂ y CO', 'Son los contaminantes menos modelados: el primer producto diario chino apareció recién en 2023.¹', 'warning'),
    ('Redes ralas', 'Fuera de estación el R² cae a 0.6 a 0.7 incluso en China¹; las macrozonas extremas de Chile son justo ese caso difícil.', 'target'),
    ('Resolución temporal', 'Predominan productos anuales o mensuales²; lo diario y horario es la frontera actual.³', 'calendar'),
]
for i, (t, d, ic) in enumerate(brechas):
    y = 1.22 + i * 0.82
    box(s, 0.5, y, 4.70, 0.72, CARD)
    icono(s, ic, 0.62, y + 0.14, 0.44, NAVY)
    text(s, 1.22, y + 0.07, 3.9, 0.22, t, size=9.3, bold=True, color=NAVY)
    text(s, 1.22, y + 0.28, 3.86, 0.42, d, size=7.5, color=INK2, leading=0.92)
box(s, 5.42, 1.22, 4.08, 3.28, DARK)
text(s, 5.62, 1.38, 3.7, 0.3, 'NUESTRO POSICIONAMIENTO', size=10, bold=True, color=GOLD)
text(s, 5.62, 1.74, 3.68, 2.65,
     [[('Aplicar la receta madura de PM₂.₅ y la frontera de los gases al territorio chileno completo: ', {'color': WHITE}),
       ('6 contaminantes, comuna × hora / día, 5 macrozonas', {'bold': True, 'color': WHITE}),
       (', con SINCA como única verdad-terreno y validación LOSO por estación.', {'color': WHITE})],
      [('El aporte es la superficie validada y reproducible; la validación epidemiológica queda '
        'explícitamente fuera del alcance, según la recomendación docente.', {'color': LIGHTBLUE})]],
     size=9.6, leading=1.04, space_after=8)
footnotes(s, ['wei23', 'vand', 'wei'], y=4.62)

# =====================================================================
# S6: implicancias y líneas base
# =====================================================================
s = S[6]
set_title(s, 'Qué implica para nuestro diseño: decisiones y líneas base', size=19)
drop_content_placeholders(s)
imps = [
    ('Predictores', 'AOD y columnas troposféricas + meteorología + uso de suelo, vías y luces nocturnas: el conjunto que la literatura repite en las cuatro familias.¹ ²', 'satellite', NAVY),
    ('Validación', 'CV espacial por estación (LOSO) y reporte por macrozona, porque el desempeño se degrada lejos de la red de monitoreo.³', 'checklist', UCBLUE),
    ('Motores', 'GWR², boosting³ y kriging cubren las tres familias estadísticas con mejor evidencia; los comparamos bajo el mismo panel.', 'science', AQUA),
]
for i, (t, d, ic, c) in enumerate(imps):
    x = 0.5 + i * 3.07
    box(s, x, 1.22, 2.93, 1.74, CARD)
    icono(s, ic, x + 0.16, 1.36, 0.44, c)
    text(s, x + 0.70, 1.38, 2.1, 0.4, t, size=10.5, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.18, 1.94, 2.6, 0.95, d, size=7.9, color=INK2, leading=0.96)
box(s, 0.5, 3.16, 9.0, 1.06, DARK)
text(s, 0.70, 3.30, 8.6, 0.3, 'LÍNEAS BASE DE DESEMPEÑO QUE FIJA LA LITERATURA', size=9.5, bold=True, color=GOLD)
text(s, 0.70, 3.62, 8.6, 0.5,
     [[('PM₂.₅ global con GWR: R² 0.81 en validación cruzada.² ', {'color': WHITE}),
       ('Gases diarios en China: R² 0.80 a 0.84 en CV y 0.61 a 0.70 fuera de estación.³ ', {'color': WHITE}),
       ('Esa es la vara realista para evaluar nuestras superficies por contaminante y macrozona.', {'color': LIGHTBLUE})]],
     size=9.3, leading=1.05)
footnotes(s, ['hoek', 'vd16', 'wei23'], y=4.44)

# =====================================================================
# S7: semana 3
# =====================================================================
s = S[7]
set_title(s, 'Semana 3: qué hicimos, qué costó y qué viene')
drop_content_placeholders(s)
colsx = [
    ('Objetivos de la semana', 'checklist', NAVY, [
        'Revisar literatura por contaminante y por método',
        'Fijar líneas base de desempeño realistas',
        'Armar el marco para la defensa de tema']),
    ('Tareas realizadas', 'download', UCBLUE, [
        'Matriz contaminante × método con referencias verificadas (DOI)',
        'Bibliografía APA del repositorio actualizada',
        'Líneas base: R² 0.8+ en CV; 0.6 a 0.7 fuera de estación']),
    ('Desafíos', 'warning', ORANGE, [
        'Literatura de SO₂ y CO escasa a nivel global',
        'Métricas no comparables entre estudios (CV muestral vs espacial)',
        'Traducir resoluciones de 1 a 10 km a comunas heterogéneas']),
    ('Próxima semana', 'calendar', AQUA, [
        'Defensa de tema (semana 4): presentación formal del problema',
        'Congelar el diseño metodológico',
        'Iniciar los paneles de modelamiento por contaminante']),
]
for i, (t, ic, c, items) in enumerate(colsx):
    x = 0.5 + i * 2.32
    box(s, x, 1.22, 2.20, 2.82, CARD)
    icono(s, ic, x + 0.14, 1.36, 0.42, c)
    text(s, x + 0.64, 1.34, 1.5, 0.5, t, size=9.5, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE, leading=0.9)
    bullets(s, x + 0.16, 1.98, 1.9, 2.0, items, size=7.6, gap=4, leading=0.95)

# =====================================================================
# S8: cierre con QR
# =====================================================================
s = S[8]
text(s, 0.55, 1.02, 5.9, 0.75, 'La revisión completa vive\nen el repositorio', size=19, bold=True, color=WHITE, leading=1.0)
bullets(s, 0.55, 1.95, 5.9, 1.7, [
    'docs/: bibliografía APA consolidada y trazable',
    'Matriz contaminante × método con cada DOI verificado',
    'Líneas base de desempeño por contaminante',
    'Todo listo para la defensa de tema de la semana 4',
], size=9.5, color=LIGHTBLUE, gap=5)
box(s, 0.55, 3.78, 5.9, 0.56, RGBColor(0x11, 0x3E, 0x8F))
s.shapes.add_picture(A('icons/github_w.png'), Inches(0.72), Inches(3.92), Inches(0.28), Inches(0.28))
text(s, 1.10, 3.78, 5.3, 0.56, 'github.com/AmaruSimonAgueroJimenez/Air-Pollution', size=11.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.95, 1.02, 2.55, 3.32, WHITE, radius=0.06)
s.shapes.add_picture(A('qr_repo.png'), Inches(7.17), Inches(1.24), Inches(2.11), Inches(2.11))
text(s, 7.05, 3.42, 2.35, 0.8, 'Escanea para ver la bibliografía, el pipeline y la documentación',
     size=8.5, color=DARK, align=PP_ALIGN.CENTER, leading=0.95)
text(s, 0.55, 4.80, 8.9, 0.3, 'Equipo 2: José J. Romero · Roberto I. Ávila · Amaru S. Agüero · Esteban A. González          ¡Gracias!',
     size=10, bold=False, color=LIGHTBLUE)

prs.save('AFG1_Video2_Presentacion_UC.pptx')
print('ok')
